#!/usr/bin/env python3

"""
Generate images from prompts using Gemini image model
and optionally assemble them into a PPTX.

Features:
- CLI arguments for all configuration
- Gemini token via argument or environment variable
- Retry + timeout logic
- Structured logging
- Clean modular design
"""

import os
import json
import logging
import argparse
from glob import glob
from io import BytesIO
from concurrent.futures import ThreadPoolExecutor, TimeoutError

from tqdm import tqdm
from PIL import Image
from google import genai

from pptx import Presentation
from pptx.util import Inches


# =========================================================
# Argument Parsing
# =========================================================

def parse_args():
    parser = argparse.ArgumentParser(
        description="Generate images using Gemini and create slides."
    )

    parser.add_argument(
        "--working-dir",
        required=True,
        help="Base working directory"
    )

    parser.add_argument(
        "--project",
        required=True,
        help="Project name"
    )

    parser.add_argument(
        "--prompts-json",
        required=True,
        help="JSON file with prompts"
    )

    parser.add_argument(
        "--model-id",
        default="gemini-3.1-flash-image-preview",
        help="Gemini model ID"
    )

    parser.add_argument(
        "--gemini-token",
        default=None,
        help="Gemini API token (if not provided, uses GEMINI_API_KEY env var)"
    )

    parser.add_argument(
        "--max-trials",
        type=int,
        default=3,
        help="Max retry attempts"
    )

    parser.add_argument(
        "--timeout-minutes",
        type=int,
        default=4,
        help="Timeout per generation"
    )

    parser.add_argument(
        "--start-prompt-index",
        type=int,
        default=0,
        help="Start index for prompts"
    )

    parser.add_argument(
        "--resolution-text",
        default="Resolution: 1920x1080. Background: white or transparent.",
        help="Extra text appended to prompts"
    )

    parser.add_argument(
        "--skip-pptx",
        action="store_true",
        help="Skip PPTX creation"
    )

    return parser.parse_args()


# =========================================================
# Token Handling
# =========================================================

def resolve_token(cli_token):
    """
    Resolve Gemini token from CLI or environment.
    """

    if cli_token:
        return cli_token

    env_token = os.environ.get("GEMINI_API_KEY")

    if not env_token:
        raise RuntimeError(
            "Gemini token not provided. "
            "Use --gemini-token or set GEMINI_API_KEY env variable."
        )

    return env_token


# =========================================================
# Directory Setup
# =========================================================

def setup_directories(working_dir, project):
    base_dir = os.path.join(working_dir, project)

    image_dir = os.path.join(base_dir, "images")
    slides_dir = os.path.join(base_dir, "slides")

    os.makedirs(image_dir, exist_ok=True)
    os.makedirs(slides_dir, exist_ok=True)

    log_file = os.path.join(base_dir, "generation.log")

    return base_dir, image_dir, slides_dir, log_file


# =========================================================
# Logging Setup
# =========================================================

def setup_logging(log_file):

    logging.basicConfig(
        filename=log_file,
        level=logging.INFO,
        format="%(asctime)s | %(levelname)s | %(message)s"
    )

    logging.info("Logging initialized")


# =========================================================
# Prompt Loading
# =========================================================

def load_prompts(json_file, resolution_text):

    with open(json_file, "r") as file:
        data = json.load(file)

    prompts = {}

    for k, v in data.items():

        idx = int(k.replace("prompt", ""))

        full_prompt = (
            v
            + " "
            + resolution_text
        )

        prompts[idx] = full_prompt.replace("\n", " ")

    sorted_keys = sorted(prompts.keys())

    return [prompts[k] for k in sorted_keys]


# =========================================================
# Image Generation
# =========================================================

def generate_image_with_retry(
        client,
        model_id,
        prompt,
        slide_index,
        image_dir,
        max_trials,
        timeout_seconds):

    def api_call():
        return client.models.generate_content(
            model=model_id,
            contents=prompt,
        )

    for trial in range(1, max_trials + 1):

        logging.info(
            f"START | Trial {trial}/{max_trials} | "
            f"Prompt: {prompt[:100]}"
        )

        try:

            with ThreadPoolExecutor(max_workers=1) as executor:

                future = executor.submit(api_call)

                response = future.result(
                    timeout=timeout_seconds
                )

            for part in response.parts:

                if part.inline_data is not None:

                    image_bytes = part.inline_data.data

                    image = Image.open(
                        BytesIO(image_bytes)
                    )

                    filename = f"{slide_index}.png"

                    filepath = os.path.join(
                        image_dir,
                        filename
                    )

                    image.save(filepath)

                    logging.info(
                        f"SUCCESS | Trial {trial} | Saved: {filepath}"
                    )

                    return True

            logging.warning(
                f"NO IMAGE DATA | Trial {trial}"
            )

        except TimeoutError:

            logging.error(
                f"TIMEOUT | Trial {trial} exceeded {timeout_seconds}s"
            )

        except Exception as e:

            logging.error(
                f"ERROR | Trial {trial} | {str(e)}"
            )

        if trial < max_trials:

            logging.info(
                f"RETRYING | Next trial {trial + 1}"
            )

        else:

            logging.error(
                f"FAILED AFTER {max_trials} TRIALS"
            )

    return False


# =========================================================
# PPTX Creation
# =========================================================

def images_to_pptx(image_folder, output_file):

    prs = Presentation()

    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    image_paths = sorted(
        glob(os.path.join(image_folder, "*.png"))
    )

    for img_path in image_paths:

        slide_layout = prs.slide_layouts[6]

        slide = prs.slides.add_slide(
            slide_layout
        )

        slide.shapes.add_picture(
            img_path,
            left=0,
            top=0,
            width=prs.slide_width,
            height=prs.slide_height
        )

    prs.save(output_file)

    print(f"Saved presentation: {output_file}")


# =========================================================
# Main
# =========================================================

def main():

    args = parse_args()

    token = resolve_token(
        args.gemini_token
    )

    (
        base_dir,
        image_dir,
        slides_dir,
        log_file
    ) = setup_directories(
        args.working_dir,
        args.project
    )

    setup_logging(log_file)

    prompts = load_prompts(
        args.prompts_json,
        args.resolution_text
    )

    timeout_seconds = (
        args.timeout_minutes * 60
    )

    client = genai.Client(
        api_key=token
    )

    total_prompts = len(prompts)

    for idx, prompt in tqdm(
            enumerate(
                prompts[
                    args.start_prompt_index:
                ],
                start=args.start_prompt_index + 1
            )):

        slide_index = str(idx).zfill(
            len(str(total_prompts))
        )

        generate_image_with_retry(
            client=client,
            model_id=args.model_id,
            prompt=prompt,
            slide_index=slide_index,
            image_dir=image_dir,
            max_trials=args.max_trials,
            timeout_seconds=timeout_seconds
        )

    if not args.skip_pptx:

        output_file = os.path.join(
            slides_dir,
            f"{args.project}.pptx"
        )

        images_to_pptx(
            image_dir,
            output_file
        )


# =========================================================

if __name__ == "__main__":
    main()
